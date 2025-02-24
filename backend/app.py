from flask import Flask, request, jsonify, make_response
import fitz  # PyMuPDF for PDF
import docx
from pptx import Presentation
import openai
import firebase_admin
from firebase_admin import credentials, storage, firestore, auth
import os
import tempfile
from flask_cors import CORS
from dotenv import load_dotenv


cred = credentials.Certificate("firebase_credentials.json")  # Update with your Firebase credentials file
firebase_admin.initialize_app(cred, {'storageBucket': 'echonotes-d6052.firebasestorage.app'})  # ✅ Correct format
db = firestore.client() 

# 🔹 Set up Flask app
app = Flask(__name__)
CORS(app)  # ✅ Enable CORS for all routes

load_dotenv()
# 🔹 OpenAI API Key
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")  # Get API key from environment
if not OPENAI_API_KEY:
    raise ValueError("Missing OpenAI API key!")

openai.api_key = OPENAI_API_KEY

# ✅ Function to extract text from files
def extract_text(file_path, file_type):
    text = ""
    if file_type == "pdf":
        doc = fitz.open(file_path)
        for page in doc:
            text += page.get_text()
    elif file_type == "docx":
        doc = docx.Document(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    elif file_type == "pptx":
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    elif file_type == "txt":
        with open(file_path, "r", encoding="utf-8") as file:
            text = file.read()
    return text.strip()

# ✅ Function to summarize text with options
def summarize_text(text, summary_type="concise"):
    if summary_type == "concise":
        max_tokens = 100  # Short summary
        prompt = f"Create a concise audiobook-style summary of the following text. Start with a compelling title, followed by a one-sentence overview. Then, present the key details structured into short, clear chapters. Keep the summary informative but streamlined, maintaining essential insights while ensuring brevity. \n{text}"
    else:  # "detailed"
        max_tokens = 300  # Longer summary
        prompt = f"Generate a comprehensive, audiobook-style summary of the following text. Begin with a title that captures the essence of the content, followed by a one-sentence overview. Then, expand into an in-depth breakdown, structured into well-defined chapters. Ensure thorough coverage of all key points, explanations, and nuances, making it engaging and detailed while maintaining clarity.  \n{text}"

    response = openai.chat.completions.create(
        model="gpt-4",
        messages=[{"role": "system", "content": prompt}],
        max_tokens=max_tokens
    )

    return response.choices[0].message.content.strip()

# ✅ Function to convert text to speech using OpenAI TTS
def convert_text_to_speech(text, filename, voice="alloy"):
    valid_voices = ["alloy", "echo", "fable", "onyx", "nova", "shimmer"]
    if voice not in valid_voices:
        voice = "alloy"  # Default voice

    response = openai.audio.speech.create(
        model="tts-1",
        input=text,
        voice=voice
    )

    # Save the audio file temporarily
    temp_audio = tempfile.NamedTemporaryFile(delete=False, suffix=".mp3")
    temp_audio.write(response.content)  # Use .content to get binary data
    temp_audio.close()
    
    # Upload to Firebase Storage
    bucket = storage.bucket()
    blob = bucket.blob(f"audiobooks/{filename}.mp3")
    blob.upload_from_filename(temp_audio.name)
    blob.make_public()
    os.remove(temp_audio.name)  # Cleanup

    return blob.public_url

# ✅ Handle Preflight (OPTIONS) Requests
@app.route("/upload", methods=["OPTIONS"])
def upload_options():
    response = make_response()
    response.headers.add("Access-Control-Allow-Origin", "*")
    response.headers.add("Access-Control-Allow-Headers", "Content-Type,Authorization")
    response.headers.add("Access-Control-Allow-Methods", "POST, OPTIONS")
    return response

@app.route("/upload", methods=["POST"])
def upload_file():
    if "file" not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file = request.files["file"]
    user_id = request.form.get("userId")
    summary_type = request.form.get("summaryType", "concise")
    voice = request.form.get("voice", "alloy")

    if not user_id:
        return jsonify({"error": "User ID required"}), 400

    file_extension = file.filename.split(".")[-1].lower()
    if file_extension not in ["pdf", "docx", "pptx", "txt"]:
        return jsonify({"error": "Invalid file type"}), 400

    # 🔹 Save temporary file
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_extension}")
    file.save(temp_file.name)

    extracted_text = extract_text(temp_file.name, file_extension)
    if not extracted_text:
        return jsonify({"error": "No text extracted"}), 400

    summary = summarize_text(extracted_text, summary_type)

    # 🔹 Generate unique filename
    temp_path = tempfile.mkstemp()[1]
    temp_filename = os.path.basename(temp_path)  # Extract only the filename
    audiobook_id = f"{user_id}_{file.filename.split('.')[0]}_{temp_filename}"
    audio_url = convert_text_to_speech(summary, audiobook_id, voice)

    # 🔹 Save audiobook metadata in Firestore
    audiobook_data = {
        "id": audiobook_id,
        "userId": user_id,
        "title": file.filename,
        "summary": summary,
        "audioUrl": audio_url,
    }

    try:
        db.collection("audiobooks").document(audiobook_id).set(audiobook_data)
        print(f"✅ Audiobook saved: {audiobook_data}")  # ✅ Debugging print statement
    except Exception as e:
        print(f"❌ Firestore Error: {e}")  # ❌ Logs Firestore errors
        return jsonify({"error": "Failed to save metadata"}), 500

    return jsonify(audiobook_data), 201


@app.route("/audiobooks", methods=["GET"])
def get_audiobooks():
    # 🔹 Get auth token from request headers
    auth_header = request.headers.get("Authorization")
    if not auth_header or not auth_header.startswith("Bearer "):
        return jsonify({"error": "Missing or invalid token"}), 401

    token = auth_header.split("Bearer ")[1]

    try:
        # 🔹 Verify Firebase token and extract user ID
        decoded_token = auth.verify_id_token(token)
        user_id = decoded_token["uid"]
    except Exception as e:
        return jsonify({"error": f"Invalid token: {str(e)}"}), 401

    # 🔹 Query Firestore for audiobooks linked to this user
    audiobooks_ref = db.collection("audiobooks").where("userId", "==", user_id)
    audiobooks = [doc.to_dict() for doc in audiobooks_ref.stream()]

    return jsonify(audiobooks)


if __name__ == "__main__":
    app.run(debug=True)